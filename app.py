import io
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt 
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from deep_translator import GoogleTranslator

st.set_page_config(page_title="5-Image PPT Generator (EN | HI | MR)", page_icon="🖼️")
st.title("🖼️ AgriSavant — 5-Image PPT Maker (English | हिंदी | मराठी)")

st.write(
    "Upload **exactly 3 images for the left section** and **2 for the right section**. "
    "The app will create **3 slides** (EN, HI, MR) with the **same images** but **translated text**."
)

# ---------- UI ----------
filename = st.text_input("Output filename", value="ags_report_multilang.pptx")
farmer_name = st.text_input("Farmer Name", value="Enter the farm name")
report_date = st.date_input("Report Date")
value_input = st.number_input("Enter the meter value", min_value=-50, max_value=100, value=50)


with st.expander("🔧 Layout controls (optional)"):
    top_in = st.slider("Top margin for images (inches)", 0.0, 5.0, 1.2, 0.1)
    side_margin_in = st.slider("Left/Right margin (inches) [kept for compatibility, not heavily used]", 0.0, 2.0, 0.5, 0.05)
    gap_in = st.slider("Gap between images (inches)", 0.0, 1.0, 0.2, 0.05)
    left_section_height_in = st.slider("Height of left section images (inches)", 0.5, 5.0, 2.0, 0.1)

# Separate uploaders for left and right images
uploaded_left = st.file_uploader(
    "Upload 3 images for the left section (side by side)", 
    type=["png", "jpg", "jpeg"], 
    accept_multiple_files=True, 
    key="left"
)
uploaded_right = st.file_uploader(
    "Upload 2 images for the right section (vertical)", 
    type=["png", "jpg", "jpeg"], 
    accept_multiple_files=True, 
    key="right"
)

# Text input for below images
below_images_text = st.text_area(
    "Text below all three left images (one point per line, will appear as separate lines in PPT)",
    value="Point 1\nPoint 2\nPoint 3\nPoint 4\nPoint 5"
)

notes_points = st.text_area(
    "Observation / Notes (one point per line, will appear as bullets in PPT)",
    value="Point 1\nPoint 2\nPoint 3\nPoint 4\nPoint 5"
)

# ---------- Translator ----------
def t(text, lang):
    try:
        return GoogleTranslator(source='auto', target=lang).translate(text)
    except Exception:
        return text

# ---------- Core PPT builder (single slide, your original logic preserved) ----------
def add_slide_with_layout(
    prs,
    images_bytes,
    farmer,
    date_str,
    lang_code,
    top=1.2,
    side_margin=0.5,
    gap=0.2,
    left_section_height=2.0,
    below_images_text="",
    notes_points="",
):
    """
    Adds ONE slide to prs using your original layout logic:
      - Title
      - 3 left images with bullets below them
      - 2 right images with Temperature/Humidity captions
      - Observation/Notes block bottom
    Translates all text to `lang_code` with googletrans.
    """

    # Translate all headings/static strings once
    title_text = t(f"Farm Report: {farmer} [{date_str}]", lang_code)
    crop_pictures_header = t("Crop Pictures", lang_code)
    weather_forecast_header = t("Weather Forecast", lang_code)
    temp_forecast_caption = t("Temperature forecast", lang_code)
    humidity_forecast_caption = t("Humidity forecast", lang_code)
    notes_header = t("OBSERVATION / NOTES:", lang_code)

    # Translate bullets
    translated_below = [t(line.strip(), lang_code) for line in below_images_text.split("\n") if line.strip()]
    translated_notes = [t(line.strip(), lang_code) for line in notes_points.split("\n") if line.strip()]

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title at the top (kept as your logic)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(198, 239, 206)  # Light green (as before)
    title_tf = title_box.text_frame
    title_run = title_tf.paragraphs[0].add_run()
    title_run.text = title_text
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    title_tf.word_wrap = True
    title_run.font.size = Pt(20)
    title_run.font.bold = True

    if len(images_bytes) != 5:
        raise ValueError("Exactly 5 images required.")

    # Divide slide into 4 equal vertical sections (kept as-is)
    slide_width_in = prs.slide_width / Inches(1)   # convert EMU -> inches
    slide_height_in = prs.slide_height / Inches(1) # convert EMU -> inches
    section_width = slide_width_in / 4

    # Section 1: 3 images side by side (use first two quarters, i.e., half slide)
    left_section_left = 0
    left_section_top = top
    left_section_width = section_width * 2  # half the slide width

    img_width_left = (left_section_width - 2 * gap) / 3
    img_height_left = left_section_height
    for i in range(3):
        img_left = left_section_left + i * (img_width_left + gap)
        slide.shapes.add_picture(
            images_bytes[i],
            Inches(img_left),
            Inches(left_section_top),
            width=Inches(img_width_left),
            height=Inches(img_height_left)
        )

    # Add a text box below all three images in the left section (with bullet points)
    textbox_top = left_section_top + img_height_left + 0.2  # 0.2 inch gap below images
    textbox_height = 1.0
    textbox = slide.shapes.add_textbox(
        Inches(left_section_left),
        Inches(textbox_top),
        Inches(left_section_width),
        Inches(textbox_height)
    )
    textbox_tf = textbox.text_frame
    textbox_tf.word_wrap = True

    if translated_below:
        textbox_tf.text = f"• {translated_below[0]}"
        textbox_tf.paragraphs[0].font.size = Pt(12)
        textbox_tf.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        for line in translated_below[1:]:
            p = textbox_tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.JUSTIFY

    # Add header above 3 left images
    header_height = 0.4
    header_box = slide.shapes.add_textbox(
        Inches(left_section_left),
        Inches(left_section_top - header_height - 0.1),  # 0.1 inch gap above images
        Inches(left_section_width),
        Inches(header_height)
    )
    header_tf = header_box.text_frame
    header_tf.text = crop_pictures_header
    header_tf.paragraphs[0].font.size = Pt(18)
    header_tf.paragraphs[0].font.bold = True
    header_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Section 2: 2 images vertically (use last two quarters, i.e., right half)
    right_section_left = section_width * 2  # Start at the middle of the slide
    right_section_top = top
    right_section_width = section_width * 2  # Use half the slide width
    right_section_height = 3  # same as your code logic

    img_height_right = (right_section_height - gap) / 2
    captions = [temp_forecast_caption, humidity_forecast_caption]
    for i in range(2):
        img_top = right_section_top + i * (img_height_right + gap)

        # Caption
        caption_box = slide.shapes.add_textbox(
            Inches(right_section_left),
            Inches(img_top),
            Inches(right_section_width),
            Inches(0.3)
        )
        caption_tf = caption_box.text_frame
        caption_tf.text = captions[i]
        caption_tf.paragraphs[0].font.size = Pt(12)
        caption_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Image below the caption
        slide.shapes.add_picture(
            images_bytes[3 + i],
            Inches(right_section_left),
            Inches(img_top + 0.3),  # 0.3 inch below the caption
            width=Inches(right_section_width),
            height=Inches(img_height_right - 0.3)  # reduce height to fit caption
        )

    # Add header above 2 right images
    right_header_height = 0.4
    right_header_box = slide.shapes.add_textbox(
        Inches(right_section_left),
        Inches(right_section_top - right_header_height - 0.1),  # 0.1 inch gap above images
        Inches(right_section_width),
        Inches(right_header_height)
    )
    right_header_tf = right_header_box.text_frame
    right_header_tf.text = weather_forecast_header
    right_header_tf.paragraphs[0].font.size = Pt(18)
    right_header_tf.paragraphs[0].font.bold = True
    right_header_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add notes/observation section centered in the bottom half (same logic)
    notes_left = section_width * 0
    notes_top = slide_height_in / 2
    notes_width = section_width * 4
    notes_height = slide_height_in / 2

    notes_box = slide.shapes.add_textbox(
        Inches(notes_left),
        Inches(notes_top),
        Inches(notes_width),
        Inches(notes_height)
    )
    notes_tf = notes_box.text_frame
    notes_tf.word_wrap = True
    notes_tf.text = notes_header
    notes_tf.paragraphs[0].font.size = Pt(16)
    notes_tf.paragraphs[0].font.bold = True
    notes_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    notes_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    for line in translated_notes:
        p = notes_tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.JUSTIFY

    # Load your logo
    logo_path = "horizontal color.png"  # Replace with your actual logo path

    # Define logo size
    logo_width_in = 2
    logo_height_in = 0.5

    # Calculate position for top-right (3rd quadrant)
    logo_left = slide_width_in - logo_width_in - 0.3  # Right-aligned with margin
    logo_top = 0.3  # Top margin

    # Add logo to slide
    slide.shapes.add_picture(
        logo_path,
        Inches(logo_left),
        Inches(logo_top),
        width=Inches(logo_width_in),
        height=Inches(logo_height_in)
    )
    
    
    # --- Add Plotly Gauge Meter in 4th quadrant (bottom-right corner) ---
    import matplotlib.pyplot as plt
    import numpy as np
    import io

    # === Gauge configuration ===
    try:
        value = float(value_input)
        if not (-50 <= value <= 100):
            st.warning("Value out of range! Please enter between -50 and 100.")
            st.stop()
    except ValueError:
        st.warning("Please enter a valid number.")
        st.stop()

    
    fig = plt.figure(figsize=(10, 6))
    ax = fig.add_subplot(projection='polar') 

    # Zones and properties
    colors = ['#4dab6d', "#f6ee54", "#ee4d55"]
    angles = [0, 1.05, 2.1]
    values = [100, 50, 0, -50]
    r_base = 2.0
    r_thickness = 0.5
    r_outer = r_base + r_thickness

    # Draw segments
    ax.bar(
        x=angles,
        width=1.05,
        height=r_thickness,
        bottom=r_base,
        linewidth=3,
        edgecolor="white",
        color=colors,
        align="edge"
    )

    # Text labels
    # ax.annotate("High Performing", xy=(0.35, r_base + 0.1), rotation=-65, color="white", fontweight="bold")
    # ax.annotate("Stable", xy=(1.55, r_base + 0.2), rotation=0, color="white", fontweight="bold")
    # ax.annotate("Needs Attention", xy=(2.55, r_base + 0.3), rotation=60, color="white", fontweight="bold")

    # Numeric markers
    marker_values = [-50, 0, 50, 100]
    marker_angles = np.interp(marker_values, [-50, 100], [0, 3.15])
    for angle_rad, val in zip(marker_angles, marker_values):
        ha = "right" if val <= 25 else "left"
        ax.annotate(
            str(val),
            xy=(angle_rad, r_outer + 0.1),
            ha=ha,
            va="center",
            fontsize=12,
            fontweight="bold",
            color="gray"
        )

    # Needle
    needle_angle = np.interp(value, [-50, 100], [0, 3.15])
    ax.annotate(
        str(value),
        xytext=(0, 0),
        xy=(needle_angle, r_base),
        arrowprops=dict(arrowstyle="wedge, tail_width=0.5", color="black", shrinkA=0),
        bbox=dict(boxstyle="circle", facecolor="black", linewidth=2.0),
        fontsize=45,
        color="white",
        ha="center"
    )

    # Boundary ring
    theta = np.linspace(0, np.pi, 300)
    r_outer_arc = np.full_like(theta, r_outer)
    r_inner_arc = np.full_like(theta, r_base)

    ax.plot(theta, r_outer_arc, color='black', linewidth=3, zorder=10)
    ax.plot(theta, r_inner_arc, color='black', linewidth=3, zorder=10)
    ax.plot([0, 0], [r_base, r_outer], color='black', linewidth=3, zorder=10)
    ax.plot([np.pi, np.pi], [r_base, r_outer], color='black', linewidth=3, zorder=10)

    # Title
    plt.title("Performance Gauge Chart", loc="center", pad=20, fontsize=35, fontweight="bold")
    ax.set_axis_off()

    # Save gauge to buffer
    gauge_buf = io.BytesIO()
    plt.savefig(gauge_buf, format='png', bbox_inches='tight', dpi=300)
    gauge_buf.seek(0)
    plt.close(fig)


    slide_width_in = prs.slide_width / Inches(1)
    slide_height_in = prs.slide_height / Inches(1)
    gauge_width_in = 2.5
    gauge_height_in = 1.5

    gauge_left = slide_width_in - gauge_width_in - 0.3
    gauge_top = slide_height_in - gauge_height_in - 0.3

    slide.shapes.add_picture(
        gauge_buf,
        Inches(gauge_left),
        Inches(gauge_top),
        width=Inches(gauge_width_in),
        height=Inches(gauge_height_in)
    )




def build_ppt_multilang(
    images_bytes,
    farmer,
    date_str,
    top=1.2,
    side_margin=0.5,
    gap=0.2,
    left_section_height=2.0,
    below_images_text="",
    notes_points=""
):
    prs = Presentation()

    # English
    add_slide_with_layout(
        prs, images_bytes, farmer, date_str,
        lang_code="en",
        top=top,
        side_margin=side_margin,
        gap=gap,
        left_section_height=left_section_height,
        below_images_text=below_images_text,
        notes_points=notes_points
    )

    # Hindi
    add_slide_with_layout(
        prs, images_bytes, farmer, date_str,
        lang_code="hi",
        top=top,
        side_margin=side_margin,
        gap=gap,
        left_section_height=left_section_height,
        below_images_text=below_images_text,
        notes_points=notes_points
    )

    # Marathi
    add_slide_with_layout(
        prs, images_bytes, farmer, date_str,
        lang_code="mr",
        top=top,
        side_margin=side_margin,
        gap=gap,
        left_section_height=left_section_height,
        below_images_text=below_images_text,
        notes_points=notes_points
    )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ---------- App Logic ----------
if uploaded_left and uploaded_right:
    if len(uploaded_left) != 3 or len(uploaded_right) != 2:
        st.error("❌ Please upload exactly **3 images for the left** and **2 images for the right**.")
    else:
        st.subheader("Preview")
        cols = st.columns(5)
        for i, f in enumerate(uploaded_left + uploaded_right):
            with cols[i]:
                st.image(f, caption=f"Image {i+1}", use_container_width=True)

        if st.button("📥 Generate Multilang PPT (EN | HI | MR)"):
            try:
                images_bytes = [io.BytesIO(f.getvalue()) for f in uploaded_left + uploaded_right]
                ppt_bytes = build_ppt_multilang(
                    images_bytes,
                    farmer=farmer_name,
                    date_str=report_date.strftime("%d-%m-%Y"),
                    top=top_in,
                    side_margin=side_margin_in,
                    gap=gap_in,
                    left_section_height=left_section_height_in,
                    below_images_text=below_images_text,
                    notes_points=notes_points
                )
                st.success("✅ PPT (3 slides) created successfully!")
                st.download_button(
                    label="⬇️ Download PPT",
                    data=ppt_bytes,
                    file_name=filename if filename.endswith(".pptx") else filename + ".pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            except Exception as e:
                st.exception(e)
else:
    st.info("📂 Upload 3 images for the left and 2 images for the right to enable the generator.")
